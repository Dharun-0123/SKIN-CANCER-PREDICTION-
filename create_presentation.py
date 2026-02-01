"""
SkinCare AI - Professional PowerPoint Presentation Generator
Creates a visually appealing presentation for formal academic presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

# Create presentation with widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme - Professional Dark Theme
PRIMARY_COLOR = RgbColor(168, 85, 247)    # Purple accent
SECONDARY_COLOR = RgbColor(6, 182, 212)   # Cyan accent
BG_DARK = RgbColor(10, 10, 15)            # Dark background
TEXT_PRIMARY = RgbColor(226, 232, 240)    # Light text
TEXT_SECONDARY = RgbColor(148, 163, 184)  # Muted text

def add_title_slide(prs, title, subtitle):
    """Add a title slide with gradient-like effect"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BG_DARK
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = SECONDARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullet_points, has_icon=False):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BG_DARK
    background.line.fill.background()
    
    # Accent bar at top
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = PRIMARY_COLOR
    accent_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content, left_title="", right_title=""):
    """Add a two-column slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BG_DARK
    background.line.fill.background()
    
    # Accent bar
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = PRIMARY_COLOR
    accent_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Left column title
    if left_title:
        left_title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(5.5), Inches(0.5))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(5.5), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(8)
    
    # Right column title
    if right_title:
        right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.5))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(5.5), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(8)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BG_DARK
    background.line.fill.background()
    
    # Accent bar
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = PRIMARY_COLOR
    accent_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(11)
    table_height = Inches(0.5 * num_rows)
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(1.1), Inches(1.5), table_width, table_height).table
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RgbColor(26, 26, 36)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_PRIMARY
            p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_title, section_number):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(18, 18, 26)
    background.line.fill.background()
    
    # Section number
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_number}"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============================================
# CREATE PRESENTATION SLIDES
# ============================================

# Slide 1: Title Slide
add_title_slide(prs, 
    "SkinCare AI",
    "AI-Powered Skin Lesion Classification System")

# Slide 2: Table of Contents
add_content_slide(prs, "Presentation Outline", [
    "Introduction & Problem Statement",
    "Objectives & Scope",
    "System Architecture",
    "Dataset Information",
    "Machine Learning Models",
    "Technology Stack",
    "Features & Functionality",
    "Results & Performance",
    "Demo & Future Enhancements",
    "Conclusion"
])

# Slide 3: Section - Introduction
add_section_slide(prs, "Introduction", 1)

# Slide 4: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "Skin cancer is one of the most common cancers globally",
    "Over 5 million cases diagnosed annually in the US alone",
    "Early detection crucial: 99% survival rate vs 27% for late-stage",
    "Limited access to dermatologists in many regions",
    "Long wait times delay critical diagnoses",
    "Professional screenings can be expensive and inaccessible",
    "Different skin conditions appear similar to untrained eyes"
])

# Slide 5: Solution Overview
add_content_slide(prs, "Our Solution: SkinCare AI", [
    "AI-powered preliminary screening tool for skin lesions",
    "Instant analysis using deep learning models",
    "Classification into 8 different skin condition categories",
    "Confidence scoring for prediction reliability",
    "Educational information about detected conditions",
    "Analysis history tracking for monitoring changes",
    "Professional PDF reports for medical consultations",
    "AI chatbot (DermaGenie) for skin health guidance"
])

# Slide 6: Section - Objectives
add_section_slide(prs, "Objectives", 2)

# Slide 7: Objectives
add_two_column_slide(prs, "Project Objectives",
    ["Develop accurate deep learning model for skin lesion classification",
     "Create user-friendly web interface for image upload",
     "Implement secure user authentication",
     "Provide educational information about conditions",
     "Generate professional reports for consultations"],
    ["Implement analytics dashboard for trends",
     "Develop AI chatbot for skin health guidance",
     "Enable comparison of multiple analyses",
     "Ensure mobile responsiveness",
     "Maintain legal compliance with disclaimers"],
    "Primary Objectives", "Secondary Objectives")

# Slide 8: Section - Architecture
add_section_slide(prs, "System Architecture", 3)

# Slide 9: Architecture Overview
add_content_slide(prs, "System Architecture", [
    "Client Layer: Browser-based access (Desktop, Mobile, Tablet)",
    "Web Server Layer: Django Web Server with Views & Templates",
    "Application Layer: User Management, Image Processing, Email Service",
    "ML/AI Layer: TensorFlow/Keras with Dual Model System",
    "Data Layer: SQLite Database + File Storage for images/models",
    "External Services: Resend (Email), Perplexity AI (Chatbot)"
])

# Slide 10: Data Flow
add_content_slide(prs, "Data Flow Pipeline", [
    "1. User uploads skin lesion image through web interface",
    "2. Image preprocessing: Resizing, Normalization",
    "3. Model selection based on user preference or auto mode",
    "4. Primary Model (EfficientNetB0) processes image",
    "5. Confidence check - if low, fallback to Secondary Model (CNN)",
    "6. Result formatting with legal compliance checks",
    "7. Display results with condition info and recommendations",
    "8. Store prediction in database and send notification"
])

# Slide 11: Section - Dataset
add_section_slide(prs, "Dataset Information", 4)

# Slide 12: Primary Dataset
add_table_slide(prs, "Primary Dataset: ISIC 2019",
    ["Class", "Full Name", "Images", "Percentage"],
    [
        ["NV", "Melanocytic Nevi", "12,875", "50.8%"],
        ["MEL", "Melanoma", "4,522", "17.9%"],
        ["BCC", "Basal Cell Carcinoma", "3,323", "13.1%"],
        ["BKL", "Benign Keratosis", "2,624", "10.4%"],
        ["AK", "Actinic Keratoses", "867", "3.4%"],
        ["SCC", "Squamous Cell Carcinoma", "628", "2.5%"],
        ["VASC", "Vascular Lesions", "253", "1.0%"],
        ["DF", "Dermatofibroma", "239", "0.9%"]
    ])

# Slide 13: Secondary Dataset
add_content_slide(prs, "Secondary Dataset: HAM10000 Subset", [
    "Base Source: HAM10000 (Human Against Machine with 10000 images)",
    "Total Images: 5,906 dermoscopic images",
    "License: CC BY-NC 4.0 (Open Source for Academic Use)",
    "8 Classes: akiec, bcc, bkl, df, mel, nv, vasc + not_skin_cancer",
    "Citation: Tschandl et al., Scientific Data, 2018",
    "Preprocessing: Resizing, Normalization, Data Augmentation",
    "Split: 70% Training / 15% Validation / 15% Testing"
])

# Slide 14: Section - ML Models
add_section_slide(prs, "Machine Learning Models", 5)

# Slide 15: Primary Model
add_table_slide(prs, "Primary Model: EfficientNetB0",
    ["Parameter", "Value"],
    [
        ["Architecture", "EfficientNetB0 (Transfer Learning)"],
        ["Input Shape", "224 × 224 × 3"],
        ["Output Classes", "8"],
        ["Total Parameters", "~5.3 Million"],
        ["Test Accuracy", "71.32%"],
        ["Training Epochs", "50"],
        ["Optimizer", "Adam (LR: 0.0001)"],
        ["Pre-trained On", "ImageNet"]
    ])

# Slide 16: Secondary Model
add_table_slide(prs, "Secondary Model: Custom CNN",
    ["Parameter", "Value"],
    [
        ["Architecture", "Custom Convolutional Neural Network"],
        ["Input Shape", "48 × 48 × 3"],
        ["Output Classes", "8"],
        ["Total Parameters", "~1.2 Million"],
        ["Test Accuracy", "94.1%"],
        ["Training Epochs", "100"],
        ["Optimizer", "Adam (LR: 0.001)"],
        ["Conv Layers", "3 Blocks (32→64→128 filters)"]
    ])

# Slide 17: Model Selection Logic
add_content_slide(prs, "Intelligent Model Selection", [
    "User can choose: EfficientNetB0, CNN, or Auto Mode",
    "Auto Mode Logic:",
    "   → First, try EfficientNetB0 (larger, more robust)",
    "   → Check confidence score of prediction",
    "   → If confidence > 50%, use EfficientNetB0 result",
    "   → If confidence ≤ 50%, fallback to CNN model",
    "Benefits: Combines strengths of both models",
    "Result: More reliable predictions across diverse images"
])

# Slide 18: Section - Technology
add_section_slide(prs, "Technology Stack", 6)

# Slide 19: Technology Stack
add_two_column_slide(prs, "Technology Stack",
    ["Python 3.10+ (Primary Language)",
     "Django 4.2.1 (Web Framework)",
     "TensorFlow 2.13.0 (Deep Learning)",
     "Keras 2.13.1 (Neural Network API)",
     "SQLite (Database)",
     "Pillow & OpenCV (Image Processing)"],
    ["HTML5, CSS3, JavaScript (Frontend)",
     "Chart.js (Data Visualization)",
     "Resend API (Email Service)",
     "Perplexity AI (Chatbot)",
     "Git (Version Control)",
     "Font Awesome (Icons)"],
    "Backend", "Frontend & Services")

# Slide 20: Section - Features
add_section_slide(prs, "Features & Functionality", 7)

# Slide 21: Core Features
add_content_slide(prs, "Core Features", [
    "Skin Lesion Analysis: AI-powered classification of skin images",
    "8 Condition Detection: Identifies 8 different skin conditions",
    "Dual Model System: EfficientNetB0 + CNN for robust predictions",
    "Confidence Scoring: Displays prediction confidence percentage",
    "User Authentication: Secure registration with email verification",
    "Password Reset: OTP-based secure password recovery",
    "Analysis History: Track and review all past analyses"
])

# Slide 22: Advanced Features
add_content_slide(prs, "Advanced Features", [
    "User Profiles: Personal info, profile pictures, statistics",
    "PDF Export: Generate professional reports for doctors",
    "Comparison Tool: Compare multiple analyses side-by-side",
    "Analytics Dashboard: Charts and trend analysis",
    "DermaGenie AI: AI chatbot for skin health guidance",
    "Admin Dashboard: System monitoring (staff only)",
    "Email Notifications: Automated alerts and updates",
    "Mobile Responsive: Works on all devices"
])

# Slide 23: Section - Results
add_section_slide(prs, "Results & Performance", 8)

# Slide 24: Performance Metrics
add_table_slide(prs, "Model Performance Summary",
    ["Metric", "EfficientNetB0", "Custom CNN"],
    [
        ["Test Accuracy", "71.32%", "94.1%"],
        ["Input Size", "224×224", "48×48"],
        ["Parameters", "5.3M", "1.2M"],
        ["Inference Time", "~200ms", "~50ms"],
        ["Training Dataset", "25,331 images", "5,906 images"],
        ["Best For", "Complex cases", "Quick screening"]
    ])

# Slide 25: Key Achievements
add_content_slide(prs, "Key Achievements", [
    "Successfully deployed dual-model AI classification system",
    "Achieved 94.1% accuracy on CNN model",
    "Implemented secure OTP-based authentication",
    "Created comprehensive analytics dashboard",
    "Developed AI-powered chatbot (DermaGenie)",
    "Built responsive UI with modern dark theme",
    "Integrated professional PDF report generation",
    "Ensured legal compliance with medical disclaimers"
])

# Slide 26: Section - Future
add_section_slide(prs, "Future Enhancements", 9)

# Slide 27: Future Enhancements
add_two_column_slide(prs, "Future Enhancements",
    ["Mobile Application (iOS/Android)",
     "Multi-language Support",
     "Enhanced Analytics",
     "RESTful API Access",
     "Batch Image Processing"],
    ["Real-time Webcam Analysis",
     "Telemedicine Integration",
     "3D Skin Mapping",
     "Predictive Risk Assessment",
     "Clinical EHR Integration"],
    "Short-term (6 months)", "Long-term (12+ months)")

# Slide 28: Conclusion
add_content_slide(prs, "Conclusion", [
    "SkinCare AI provides accessible AI-powered skin cancer screening",
    "Dual-model architecture ensures robust and reliable predictions",
    "User-friendly interface makes technology accessible to everyone",
    "Comprehensive features: History, Analytics, Reports, AI Chat",
    "Encourages early detection and professional medical consultation",
    "Built with modern technologies and best practices",
    "Potential to save lives through early skin cancer detection"
])

# Slide 29: Disclaimer
add_content_slide(prs, "Medical Disclaimer", [
    "SkinCare AI is for EDUCATIONAL purposes only",
    "Does NOT provide medical diagnoses or treatment advice",
    "Results are AI predictions based on image analysis",
    "Should NOT replace professional medical evaluation",
    "Always consult a qualified dermatologist for:",
    "   → Proper medical evaluation",
    "   → Accurate diagnosis",
    "   → Appropriate treatment recommendations"
])

# Slide 30: Thank You
add_title_slide(prs, 
    "Thank You",
    "Questions & Discussion")

# Slide 31: References
add_content_slide(prs, "References", [
    "ISIC Archive - International Skin Imaging Collaboration",
    "HAM10000 Dataset - Tschandl et al., Scientific Data, 2018",
    "EfficientNet: Rethinking Model Scaling for CNNs (Tan & Le, 2019)",
    "Django Documentation - docs.djangoproject.com",
    "TensorFlow Documentation - tensorflow.org",
    "Skin Cancer Foundation - skincancer.org",
    "Kaggle - Skin Cancer MNIST: HAM10000 Dataset"
])

# Save the presentation
prs.save('SkinCare_AI_Presentation.pptx')
print("✅ Presentation created successfully: SkinCare_AI_Presentation.pptx")
print("📊 Total slides: 31")
print("\nSlide Structure:")
print("  1. Title Slide")
print("  2. Table of Contents")
print("  3-5. Introduction & Problem Statement")
print("  6-7. Objectives")
print("  8-10. System Architecture")
print("  11-13. Dataset Information")
print("  14-17. Machine Learning Models")
print("  18-19. Technology Stack")
print("  20-22. Features & Functionality")
print("  23-25. Results & Performance")
print("  26-27. Future Enhancements")
print("  28-29. Conclusion & Disclaimer")
print("  30-31. Thank You & References")
